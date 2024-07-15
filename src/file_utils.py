import PyPDF2  # Library for reading PDF files
import docx  # Library for reading DOCX files
from pptx import Presentation  # Library for reading PPTX files
import pandas as pd  # Library for data manipulation and analysis
import json  # Library for parsing JSON
import ast  # Library for parsing Python abstract syntax trees
from image_utils import process_image_with_neva  # Function for processing images with Neva

def parse_pdf(file_path):
    """
    Parse a PDF file and extract text from it.
    """
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)  # Initialize PDF reader
            text = ""
            for page in reader.pages:  # Iterate through all pages
                text += page.extract_text() + "\n"  # Extract text from each page
        return text
    except Exception as e:
        return f"Error parsing PDF file: {str(e)}"  # Return error message if any exception occurs

def parse_docx(file_path):
    """
    Parse a DOCX file and extract text from it.
    """
    try:
        doc = docx.Document(file_path)  # Open DOCX document
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])  # Extract text from each paragraph
    except Exception as e:
        return f"Error parsing DOCX file: {str(e)}"  # Return error message if any exception occurs

def parse_pptx(file_path):
    """
    Parse a PPTX file and extract text from it.
    """
    try:
        prs = Presentation(file_path)  # Open PPTX presentation
        text = ""
        for slide in prs.slides:  # Iterate through all slides
            for shape in slide.shapes:  # Iterate through all shapes
                if hasattr(shape, 'text'):  # Check if shape has text attribute
                    text += shape.text + "\n"  # Extract text from shape
        return text
    except Exception as e:
        return f"Error parsing PPTX file: {str(e)}"  # Return error message if any exception occurs

def parse_csv(file_path):
    """
    Parse a CSV file and convert it to a string.
    """
    try:
        df = pd.read_csv(file_path)  # Read CSV file into DataFrame
        return df.to_string()  # Convert DataFrame to string
    except Exception as e:
        return f"Error parsing CSV file: {str(e)}"  # Return error message if any exception occurs

def parse_xlsx(file_path):
    """
    Parse an XLSX file and convert it to a string.
    """
    try:
        df = pd.read_excel(file_path)  # Read Excel file into DataFrame
        return df.to_string()  # Convert DataFrame to string
    except Exception as e:
        return f"Error parsing XLSX file: {str(e)}"  # Return error message if any exception occurs

def parse_txt(file_path):
    """
    Parse a TXT file and read its content.
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()  # Read and return the content of the file
    except Exception as e:
        return f"Error parsing TXT file: {str(e)}"  # Return error message if any exception occurs

def parse_python(file_path):
    """
    Parse a Python file and convert its abstract syntax tree (AST) to a string.
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()  # Read file content
        tree = ast.parse(content)  # Parse content into an AST
        return ast.dump(tree, annotate_fields=False)  # Convert AST to string
    except SyntaxError as e:
        return f"Error parsing Python file: {str(e)}\n\nRaw content:\n{content}"  # Return error message for syntax errors
    except Exception as e:
        return f"Error reading or parsing Python file: {str(e)}"  # Return error message for other exceptions

def parse_json(file_path):
    """
    Parse a JSON file and format its content with indentation.
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)  # Load JSON content
            return json.dumps(data, indent=2)  # Return pretty-printed JSON string
    except json.JSONDecodeError as e:
        return f"Error parsing JSON file: {str(e)}\n\nRaw content:\n{file.read()}"  # Return error message for JSON decoding errors
    except Exception as e:
        return f"Error reading or parsing JSON file: {str(e)}"  # Return error message for other exceptions

def process_file(file_path):
    """
    Process a file based on its extension and extract relevant content.
    """
    file_extension = file_path.split('.')[-1].lower()  # Get the file extension
    content = ""
    if file_extension == 'pdf':
        content = parse_pdf(file_path)  # Process PDF file
    elif file_extension == 'docx':
        content = parse_docx(file_path)  # Process DOCX file
    elif file_extension == 'pptx':
        content = parse_pptx(file_path)  # Process PPTX file
    elif file_extension == 'csv':
        content = parse_csv(file_path)  # Process CSV file
    elif file_extension == 'xlsx':
        content = parse_xlsx(file_path)  # Process XLSX file
    elif file_extension == 'txt':
        content = parse_txt(file_path)  # Process TXT file
    elif file_extension in ['png', 'jpg', 'jpeg']:
        content = process_image_with_neva(file_path)  # Process image file
    elif file_extension == 'py':
        content = parse_python(file_path)  # Process Python file
    elif file_extension == 'json':
        content = parse_json(file_path)  # Process JSON file
    else:
        content = f"Unsupported file type: {file_extension}"  # Return error message for unsupported file types
    
    return str(content)  # Ensure the return value is always a string
